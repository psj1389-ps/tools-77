import os, math
import pdfplumber
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from collections import defaultdict

# OCR 준비
try:
    import pytesseract
    _t_paths = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"
    ]
    for _p in _t_paths:
        if os.path.exists(_p):
            pytesseract.pytesseract.tesseract_cmd = _p
            break
    OCR_OK = True
except ImportError:
    pytesseract = None
    OCR_OK = False

def log(msg):
    print(msg, flush=True)

def looks_garbled(text: str) -> bool:
    """텍스트 깨짐 감지"""
    if not text:
        return True
    s = text[:1000]
    hangul = sum(1 for c in s if 0xAC00 <= ord(c) <= 0xD7A3)
    letters = sum(1 for c in s if c.isalpha() or 0xAC00 <= ord(c) <= 0xD7A3)
    if letters == 0:
        return True
    junk = sum(1 for c in s if (0 <= ord(c) < 32) or (127 <= ord(c) < 160))
    return (hangul / letters) < 0.05 or (junk / len(s)) > 0.25

def extract_text_pdf(pdf_path, page_index):
    """PDF에서 텍스트 추출 (pdfplumber 사용)"""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            if page_index >= len(pdf.pages):
                return ""
            page = pdf.pages[page_index]
            
            # words 기반 라인 재구성
            words = page.extract_words(use_text_flow=True) or []
            if words:
                line_map = defaultdict(list)
                for w in words:
                    mid = (w["top"] + w["bottom"]) / 2
                    key = round(mid / 3)
                    line_map[key].append(w)
                
                lines = []
                for k in sorted(line_map.keys()):
                    ws = sorted(line_map[k], key=lambda x: x["x0"])
                    lines.append(" ".join(x["text"] for x in ws))
                
                if lines:
                    return "\n".join(lines)
            
            # fallback
            return page.extract_text() or ""
    except Exception as e:
        log(f"[extract_text_pdf] 오류: {e}")
        return ""

def extract_text_ocr(pdf_path, page_index, dpi=200):
    """OCR을 사용한 텍스트 추출"""
    if not (OCR_OK and pytesseract):
        return ""
    try:
        imgs = convert_from_path(pdf_path, dpi=dpi,
                               first_page=page_index + 1,
                               last_page=page_index + 1)
        if not imgs:
            return ""
        return pytesseract.image_to_string(imgs[0], lang="kor+eng").strip()
    except Exception as e:
        log(f"[OCR] 실패 p{page_index}: {e}")
        return ""

def get_clean_text(pdf_path, page_index):
    """깨끗한 텍스트 추출 (자동 OCR 백업)"""
    base = extract_text_pdf(pdf_path, page_index)
    if looks_garbled(base):
        log(f"[p{page_index}] 텍스트 깨짐 감지 → OCR")
        ocr = extract_text_ocr(pdf_path, page_index)
        if ocr and not looks_garbled(ocr):
            log(f"[p{page_index}] OCR 사용")
            return ocr
        # OCR도 만족 못하면 base 혹은 ocr 중 길이 긴 것 반환
        return ocr if len(ocr) > len(base) else base
    return base

def add_table_chunk(slide, numbered_lines, top_in=4.5, height_in=3.0):
    """
    슬라이드에 번호가 매겨진 텍스트 표 추가
    numbered_lines: [(번호, 텍스트), ...]
    """
    if not numbered_lines:
        return
    
    rows = len(numbered_lines) + 1  # 헤더 포함
    cols = 2
    
    try:
        table_shape = slide.shapes.add_table(rows, cols,
                                            Inches(0), Inches(top_in),
                                            Inches(10), Inches(height_in))
        tbl = table_shape.table
        
        # 헤더 설정
        tbl.cell(0, 0).text = "No"
        tbl.cell(0, 1).text = "Text"
        
        # 헤더 스타일
        for c in range(cols):
            cell = tbl.cell(0, c)
            # 배경색 설정
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(220, 220, 220)  # 연한 회색
            
            # 폰트 설정
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)
                    run.font.name = "맑은 고딕"
                    run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 내용 채우기
        for i, (num, line) in enumerate(numbered_lines, start=1):
            if i < rows:  # 행 수 초과 방지
                tbl.cell(i, 0).text = str(num)
                tbl.cell(i, 1).text = line[:500]  # 텍스트 길이 제한
                
                # 내용 스타일
                for col in range(cols):
                    cell = tbl.cell(i, col)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            run.font.name = "맑은 고딕"
                            run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 테이블 스타일 개선
        tbl.first_row = True
        
    except Exception as e:
        log(f"표 생성 오류: {e}")
        # 표 생성 실패 시 텍스트박스로 대체
        add_text_fallback(slide, numbered_lines, top_in, height_in)

def add_text_fallback(slide, numbered_lines, top_in, height_in):
    """표 생성 실패 시 텍스트박스 대체"""
    try:
        textbox = slide.shapes.add_textbox(Inches(0), Inches(top_in),
                                         Inches(10), Inches(height_in))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # 배경 설정
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 249, 250)
        
        # 텍스트 내용
        lines = [f"{num}. {text}" for num, text in numbered_lines]
        text_frame.text = "\n".join(lines)
        
        # 폰트 설정
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)
                run.font.name = "맑은 고딕"
                run.font.color.rgb = RGBColor(0, 0, 0)
                
    except Exception as e:
        log(f"텍스트박스 생성도 실패: {e}")

def add_page_image(slide, img_path):
    """슬라이드에 페이지 이미지 추가"""
    try:
        slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                               Inches(10), Inches(4.5))
    except Exception as e:
        log(f"이미지 추가 실패: {e}")
        # 이미지 추가 실패 시 플레이스홀더
        textbox = slide.shapes.add_textbox(Inches(0), Inches(0),
                                         Inches(10), Inches(4.5))
        textbox.text_frame.text = "이미지 로드 실패"

def page_to_images(pdf_path, page_index, dpi=160):
    """PDF 페이지를 이미지로 변환"""
    try:
        imgs = convert_from_path(pdf_path, dpi=dpi,
                               first_page=page_index+1,
                               last_page=page_index+1)
        return imgs
    except Exception as e:
        log(f"페이지 이미지 변환 실패: {e}")
        return []

def split_lines(lines, max_lines_per_slide):
    """텍스트 라인을 슬라이드별로 분할"""
    for i in range(0, len(lines), max_lines_per_slide):
        yield lines[i:i+max_lines_per_slide]

def convert_pdf(pdf_path, output_path,
               dpi_image=160,
               max_lines_per_slide=20,  # 표 높이 고려하여 줄임
               table_height_in=3.0):
    """PDF를 PPTX로 변환"""
    try:
        prs = Presentation()
        blank = prs.slide_layouts[6]  # 빈 레이아웃
        
        # PDF 페이지 수 확인
        with pdfplumber.open(pdf_path) as pdf:
            total_pages = len(pdf.pages)
        
        log(f"총 {total_pages}페이지 변환 시작")
        
        for p in range(total_pages):
            log(f"=== 페이지 {p+1}/{total_pages} 처리 시작 ===")
            
            # 텍스트 추출
            raw_text = get_clean_text(pdf_path, p)
            
            # 라인 전처리
            lines = [l.strip() for l in raw_text.splitlines()]
            lines = [l for l in lines if l and len(l) > 1]  # 의미있는 라인만
            
            if not lines:
                lines = ["(텍스트 없음)"]
            
            # 번호 매기기
            numbered = list(enumerate(lines, start=1))
            chunks = list(split_lines(numbered, max_lines_per_slide))
            
            log(f"[p{p+1}] 원문 라인수: {len(lines)}, 슬라이드 분할: {len(chunks)}")
            
            # 슬라이드 생성
            for ci, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(blank)
                
                # 첫 번째 청크에만 이미지 추가
                if ci == 0:
                    try:
                        imgs = page_to_images(pdf_path, p, dpi=dpi_image)
                        if imgs:
                            tmp = f"__temp_p{p}_{os.getpid()}.png"
                            imgs[0].save(tmp, "PNG")
                            add_page_image(slide, tmp)
                            
                            # 임시 파일 정리
                            try:
                                os.remove(tmp)
                            except:
                                pass
                    except Exception as e:
                        log(f"[p{p+1}] 이미지 변환 실패: {e}")
                else:
                    # 연속 슬라이드 표시
                    try:
                        box = slide.shapes.add_textbox(Inches(0), Inches(0),
                                                      Inches(10), Inches(0.6))
                        tf = box.text_frame
                        tf.text = f"페이지 {p+1} (계속 - {ci+1}/{len(chunks)})"
                        
                        # 스타일 설정
                        for paragraph in tf.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
                                run.font.size = Pt(14)
                                run.font.name = "맑은 고딕"
                                run.font.color.rgb = RGBColor(0, 100, 200)
                    except Exception as e:
                        log(f"헤더 추가 실패: {e}")
                
                # 표 추가
                add_table_chunk(slide, chunk, top_in=4.5, height_in=table_height_in)
        
        # PPTX 저장
        prs.save(output_path)
        log(f"✅ 변환 완료: {output_path}")
        return True
        
    except Exception as e:
        log(f"❌ 변환 실패: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    """메인 실행 함수"""
    # 테스트용 설정
    test_pdf = "sample.pdf"
    out_pptx = "converted_result.pptx"
    
    # uploads 폴더에서 PDF 파일 찾기
    upload_dir = "uploads/pdf"
    if os.path.exists(upload_dir):
        pdf_files = [f for f in os.listdir(upload_dir) if f.lower().endswith('.pdf')]
        if pdf_files:
            test_pdf = os.path.join(upload_dir, pdf_files[0])
            log(f"발견된 PDF 파일: {test_pdf}")
    
    if os.path.exists(test_pdf):
        log(f"변환 시작: {test_pdf} → {out_pptx}")
        success = convert_pdf(test_pdf, out_pptx,
                            dpi_image=180,  # 품질 향상
                            max_lines_per_slide=18,
                            table_height_in=3.2)
        if success:
            log(f"🎉 변환 성공! 결과: {out_pptx}")
        else:
            log("💥 변환 실패")
    else:
        log(f"❌ PDF 파일을 찾을 수 없습니다: {test_pdf}")
        log("💡 uploads/pdf/ 폴더에 PDF 파일을 넣거나 sample.pdf를 현재 폴더에 두세요.")

if __name__ == "__main__":
    main()