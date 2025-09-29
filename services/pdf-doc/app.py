from flask import Flask, request, render_template, send_file, flash, redirect, url_for, jsonify
from flask_cors import CORS
import os
import tempfile
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image
import json
from dotenv import load_dotenv
from docx import Document
from docx.shared import Pt, Inches as DocxInches
import subprocess
import platform
import pytesseract
import cv2
import numpy as np
import fitz  # PyMuPDF
import re
from typing import List, Tuple, Dict, Any
from pdf2docx import Converter
# Adobe PDF Services SDK 임포트 및 설정
try:
    # 올바른 Adobe PDF Services SDK import 구문
    from adobe.pdfservices.operation.auth.service_principal_credentials import ServicePrincipalCredentials
    from adobe.pdfservices.operation.exception.exceptions import ServiceApiException, ServiceUsageException, SdkException
    from adobe.pdfservices.operation.io.cloud_asset import CloudAsset
    from adobe.pdfservices.operation.io.stream_asset import StreamAsset
    from adobe.pdfservices.operation.pdf_services import PDFServices
    from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
    from adobe.pdfservices.operation.pdfjobs.jobs.export_pdf_job import ExportPDFJob
    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_params import ExportPDFParams
    from adobe.pdfservices.operation.pdfjobs.params.export_pdf.export_pdf_target_format import ExportPDFTargetFormat
    from adobe.pdfservices.operation.pdfjobs.result.export_pdf_result import ExportPDFResult
    import traceback
    
    adobe_available = True
    ADOBE_SDK_AVAILABLE = True
    print("Adobe PDF Services SDK가 성공적으로 로드되었습니다.")
except ImportError as e:
    print(f"Adobe PDF Services SDK를 가져올 수 없습니다: {e}")
    print("Adobe SDK 없이 계속 진행합니다.")
    adobe_available = False
    ADOBE_SDK_AVAILABLE = False

# Adobe SDK 가용성은 import 성공/실패에 따라 자동 결정됨

# 환경 변수 로드
load_dotenv()

# 환경변수 기반 설정
MAX_FILE_SIZE_MB = int(os.environ.get('MAX_FILE_SIZE_MB', '100'))
ENABLE_DEBUG_LOGS = os.getenv('ENABLE_DEBUG_LOGS', 'true').lower() == 'true'
CONVERSION_TIMEOUT = int(os.environ.get('CONVERSION_TIMEOUT_SECONDS', '300'))
TEMP_FILE_CLEANUP = os.environ.get('TEMP_FILE_CLEANUP', 'true').lower() == 'true'

app = Flask(__name__)
CORS(app, origins=["https://tools-77.vercel.app", "http://localhost:3000"])  # CORS 설정 추가
app.secret_key = os.environ.get("SECRET_KEY", "dev-secret-change-me")
max_mb = int(os.environ.get("MAX_CONTENT_LENGTH_MB", "100"))
app.config["MAX_CONTENT_LENGTH"] = max_mb * 1024 * 1024

# 헬스체크
@app.route("/health")
def health():
    return "ok", 200

@app.route("/env-check")
def env_check():
    """환경변수 설정 상태 확인 (디버깅용)"""
    # Adobe API 사용 가능성 확인
    adobe_ready = is_adobe_api_available()
    
    env_status = {
        "adobe_sdk_available": ADOBE_SDK_AVAILABLE,
        "adobe_api_ready": adobe_ready,
        "fallback_mode": not adobe_ready,
        "conversion_method": "Adobe API" if adobe_ready else "pdf2docx + OCR",
        "environment_variables": {
            "ADOBE_CLIENT_ID": "설정됨" if os.getenv('ADOBE_CLIENT_ID') else "미설정",
            "ADOBE_CLIENT_SECRET": "설정됨" if os.getenv('ADOBE_CLIENT_SECRET') else "미설정",
            "ADOBE_ORGANIZATION_ID": "설정됨" if os.getenv('ADOBE_ORGANIZATION_ID') else "미설정",
            "ADOBE_ACCOUNT_ID": "설정됨" if os.getenv('ADOBE_ACCOUNT_ID') else "미설정",
            "ADOBE_TECHNICAL_ACCOUNT_EMAIL": "설정됨" if os.getenv('ADOBE_TECHNICAL_ACCOUNT_EMAIL') else "미설정"
        },
        "app_settings": {
            "max_file_size_mb": MAX_FILE_SIZE_MB,
            "debug_logs_enabled": ENABLE_DEBUG_LOGS,
            "conversion_timeout_seconds": CONVERSION_TIMEOUT,
            "temp_file_cleanup": TEMP_FILE_CLEANUP
        },
        "config_values": {
            "client_id_length": len(os.getenv('ADOBE_CLIENT_ID', '')),
            "client_secret_length": len(os.getenv('ADOBE_CLIENT_SECRET', '')),
            "organization_id_length": len(os.getenv('ADOBE_ORGANIZATION_ID', ''))
        },
        "authentication_method": "OAuth Server-to-Server",
        "service_status": {
            "pdf2docx_available": True,  # 항상 사용 가능
            "ocr_available": True,       # 항상 사용 가능
            "image_conversion_available": True  # 항상 사용 가능
        },
        "recommendations": get_environment_recommendations(adobe_ready)
    }
    return jsonify(env_status)

def get_environment_recommendations(adobe_ready):
    """환경 설정에 대한 권장사항 제공"""
    recommendations = []
    
    if not adobe_ready:
        if not ADOBE_SDK_AVAILABLE:
            recommendations.append("Adobe PDF Services SDK가 설치되지 않았습니다. 고급 기능을 위해 설치를 고려해보세요.")
        else:
            recommendations.append("Adobe API 환경변수를 설정하면 더 나은 PDF 변환 품질을 얻을 수 있습니다.")
            recommendations.append("로컬 개발: .env 파일에 Adobe API 키를 추가하세요.")
            recommendations.append("배포 환경: 환경변수로 Adobe API 키를 설정하세요.")
    
    if not recommendations:
        recommendations.append("모든 설정이 올바르게 구성되었습니다! (OAuth Server-to-Server 인증 사용)")
    
    return recommendations

# Adobe PDF Services API 설정 - 환경변수에서 로드
# Adobe private key path function removed for OAuth Server-to-Server authentication

ADOBE_CONFIG = {
    "client_credentials": {
        "client_id": os.getenv('ADOBE_CLIENT_ID', ''),
        "client_secret": os.getenv('ADOBE_CLIENT_SECRET', '')
    },
    "service_principal_credentials": {
        "organization_id": os.getenv('ADOBE_ORGANIZATION_ID', ''),
        "account_id": os.getenv('ADOBE_ACCOUNT_ID', ''),
        "technical_account_email": os.getenv('ADOBE_TECHNICAL_ACCOUNT_EMAIL', ''),
        "access_token": ''  # 동적으로 생성되어야 함
    }
}

# Adobe API 사용 가능성 확인
def is_adobe_api_available():
    """Adobe API 사용 가능 여부 확인"""
    if not ADOBE_SDK_AVAILABLE:
        print("❌ Adobe SDK가 설치되지 않음 - Adobe API 사용 불가")
        return False
    
    client_id = ADOBE_CONFIG["client_credentials"]["client_id"]
    client_secret = ADOBE_CONFIG["client_credentials"]["client_secret"]
    organization_id = ADOBE_CONFIG["service_principal_credentials"]["organization_id"]
    account_id = ADOBE_CONFIG["service_principal_credentials"]["account_id"]
    
    # 모든 필수 자격증명 확인
    has_credentials = bool(client_id and client_secret and organization_id and account_id)
    
    print(f"🔍 Adobe API 자격증명 상태 확인:")
    print(f"  - ADOBE_CLIENT_ID: {'✅' if client_id else '❌'} {'(' + client_id[:8] + '...)' if client_id else '(누락)'}")
    print(f"  - ADOBE_CLIENT_SECRET: {'✅' if client_secret else '❌'} {'(설정됨)' if client_secret else '(누락)'}")
    print(f"  - ADOBE_ORGANIZATION_ID: {'✅' if organization_id else '❌'} {'(' + organization_id[:8] + '...)' if organization_id else '(누락)'}")
    print(f"  - ADOBE_ACCOUNT_ID: {'✅' if account_id else '❌'} {'(' + account_id[:8] + '...)' if account_id else '(누락)'}")
    
    if not has_credentials:
        missing = []
        if not client_id: missing.append('ADOBE_CLIENT_ID')
        if not client_secret: missing.append('ADOBE_CLIENT_SECRET')
        if not organization_id: missing.append('ADOBE_ORGANIZATION_ID')
        if not account_id: missing.append('ADOBE_ACCOUNT_ID')
        print(f"❌ Adobe API 사용 불가 - 누락된 환경변수: {', '.join(missing)}")
        return False
    
    print("✅ Adobe API 자격증명 완료 - OAuth Server-to-Server 인증 준비됨")
    return True

# Adobe SDK 상태 확인 및 초기화
print(f"Adobe SDK 가용성: {ADOBE_SDK_AVAILABLE}")
adobe_api_ready = is_adobe_api_available()

if adobe_api_ready:
    client_id = ADOBE_CONFIG['client_credentials']['client_id']
    print(f"✅ Adobe API 준비 완료 (OAuth Server-to-Server): {client_id[:8]}...")
else:
    print("⚠️ Adobe API 사용 불가 - fallback 모드로 작동합니다.")
    if not ADOBE_SDK_AVAILABLE:
        print("  - Adobe SDK가 설치되지 않음")
    else:
        print("  - Adobe API 환경변수가 설정되지 않음")
    print("  - pdf2docx 및 OCR 방법을 사용합니다.")

# Adobe API 가용성을 전역 변수로 설정
adobe_available = adobe_api_ready



UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'jpg', 'jpeg', 'png', 'gif', 'bmp'}

# 폴더 생성
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs('debug_output', exist_ok=True)  # 디버깅용 폴더

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# 디버깅용 중간 결과물 저장 함수들
def save_debug_text(text, filename_prefix):
    """추출된 텍스트를 디버깅용 .txt 파일로 저장"""
    try:
        debug_file = os.path.join('debug_output', f'{filename_prefix}_extracted_text.txt')
        with open(debug_file, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"디버깅: 텍스트가 {debug_file}에 저장되었습니다. (길이: {len(text)}자)")
        return debug_file
    except Exception as e:
        print(f"디버깅 텍스트 저장 오류: {e}")
        return None

def save_debug_image(image, filename_prefix, page_num):
    """변환된 이미지를 디버깅용 .png 파일로 저장"""
    try:
        debug_file = os.path.join('debug_output', f'{filename_prefix}_page_{page_num}.png')
        image.save(debug_file, 'PNG')
        print(f"디버깅: 이미지가 {debug_file}에 저장되었습니다.")
        return debug_file
    except Exception as e:
        print(f"디버깅 이미지 저장 오류: {e}")
        return None

def pdf_to_docx_with_pdf2docx(pdf_path, output_path):
    """pdf2docx 라이브러리를 사용한 PDF → DOCX 변환"""
    try:
        print("pdf2docx 라이브러리를 사용하여 변환 중...")
        
        # Converter 객체 생성
        cv = Converter(pdf_path)
        
        # 변환 실행
        cv.convert(output_path, start=0, end=None)
        
        # 객체 닫기
        cv.close()
        
        print(f"pdf2docx 변환 완료: {output_path}")
        return True
        
    except Exception as e:
        print(f"pdf2docx 변환 실패: {e}")
        return False

def ocr_image_to_blocks(pil_image):
    """이미지에서 단어 단위 텍스트와 위치(좌표)를 추출"""
    try:
        # OCR 가용성 확인
        try:
            import pytesseract
            # Tesseract 경로 자동 감지 (Render 환경 대응)
            if os.path.exists('/usr/bin/tesseract'):
                pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
        except ImportError:
            print("pytesseract를 사용할 수 없습니다.")
            return []
        
        img = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray = cv2.medianBlur(gray, 3)
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        gray = clahe.apply(gray)

        try:
            config = r"--oem 3 --psm 6 -l kor+eng"
            data = pytesseract.image_to_data(gray, config=config,
                                             output_type=pytesseract.Output.DICT)
        except Exception as ocr_error:
            print(f"한국어+영어 OCR 실패: {ocr_error}")
            # Fallback: 영어만으로 재시도
            try:
                config = r"--oem 3 --psm 6 -l eng"
                data = pytesseract.image_to_data(gray, config=config,
                                                 output_type=pytesseract.Output.DICT)
                print("영어 OCR로 fallback 성공")
            except Exception:
                print("OCR 완전 실패")
                return []
        
        blocks = []
        n = len(data["text"])
        for i in range(n):
            text = data["text"][i].strip()
            conf_val = data["conf"][i]
            if isinstance(conf_val, (int, float)):
                conf = int(conf_val)
            elif isinstance(conf_val, str) and conf_val.replace('.', '').replace('-', '').isdigit():
                conf = int(float(conf_val))
            else:
                conf = 0
            
            if conf > 30 and len(text) > 0:
                x, y, w, h = data["left"][i], data["top"][i], data["width"][i], data["height"][i]
                blocks.append({
                    "text": text,
                    "bbox": (x, y, x + w, y + h),
                    "confidence": conf
                })
        return blocks
    except Exception as e:
        print(f"OCR 처리 중 오류: {e}")
        return []

def clean_special_characters(text: str) -> str:
    """특수 문자 처리 개선 - PDF에서 잘못 추출되는 문자들을 올바르게 복구"""
    if not text:
        return text
    
    # 일반적인 PDF 추출 오류 수정
    replacements = {
        '\uf0b7': '•',  # 불릿 포인트
        '\uf0a7': '§',  # 섹션 기호
        '\uf0e0': '→',  # 화살표
        '\u2022': '•',  # 불릿 포인트
        '\u201C': '"',  # 왼쪽 큰따옴표
        '\u201D': '"',  # 오른쪽 큰따옴표
        '\u2018': "'",  # 왼쪽 작은따옴표
        '\u2019': "'",  # 오른쪽 작은따옴표
        '\u2013': '–',  # en dash
        '\u2014': '—',  # em dash
        '\u00A0': ' ',  # 줄바꿈 없는 공백
        '\u200B': '',   # 폭이 0인 공백
        '\uFEFF': '',   # 바이트 순서 표시
    }
    
    # 특수 문자 변환
    for old, new in replacements.items():
        text = text.replace(old, new)
    
    # 연속된 공백 정리
    text = re.sub(r'[\s\t\n\r]+', ' ', text)
    
    # 제로 폭 문자 제거
    text = re.sub(r'[\u200B-\u200D\uFEFF]', '', text)
    
    return text.strip()

def analyze_pdf_orientation(pdf_path: str) -> Dict[str, Any]:
    """PDF 페이지 크기를 분석하여 문서 방향 감지 (pdfplumber 사용)"""
    try:
        import pdfplumber
        page_orientations = []
        
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                width = page.width
                height = page.height
                
                # 가로/세로 방향 판단
                if width > height:
                    orientation = 'landscape'  # 가로형
                else:
                    orientation = 'portrait'   # 세로형
                
                page_orientations.append({
                    'page': page_num,
                    'width': width,
                    'height': height,
                    'orientation': orientation,
                    'aspect_ratio': width / height
                })
        
        # 전체 문서의 주요 방향 결정
        landscape_count = sum(1 for p in page_orientations if p['orientation'] == 'landscape')
        portrait_count = len(page_orientations) - landscape_count
        
        primary_orientation = 'landscape' if landscape_count > portrait_count else 'portrait'
        
        return {
            'page_orientations': page_orientations,
            'primary_orientation': primary_orientation,
            'landscape_pages': landscape_count,
            'portrait_pages': portrait_count,
            'total_pages': len(page_orientations)
        }
        
    except Exception as e:
        print(f"PDF 방향 분석 중 오류: {e}")
        return {
            'page_orientations': [],
            'primary_orientation': 'portrait',
            'landscape_pages': 0,
            'portrait_pages': 0,
            'total_pages': 0
        }

def adobe_pdf_to_docx(input_path: str, output_path: str):
    print(">>> [DEBUG] adobe_pdf_to_docx: start", flush=True)
    info = {}
    try:
        # 파일 크기 및 기본 정보 확인
        file_size = os.path.getsize(input_path)
        print(f">>> [DEBUG] 파일 크기: {file_size} bytes ({file_size/1024/1024:.2f} MB)", flush=True)
        
        # Adobe API 파일 크기 제한 확인 (100MB)
        if file_size > 100 * 1024 * 1024:
            print(">>> [DEBUG] 파일이 Adobe API 제한(100MB)을 초과함", flush=True)
            info["error"] = "FILE_TOO_LARGE_FOR_ADOBE"
            return False, info
        
        # PDF 파일 유효성 검사
        with open(input_path, "rb") as f:
            header = f.read(8)
            if not header.startswith(b'%PDF-'):
                print(">>> [DEBUG] 유효하지 않은 PDF 헤더", flush=True)
                info["error"] = "INVALID_PDF_HEADER"
                return False, info
        
        creds = ServicePrincipalCredentials(
            client_id=os.environ["ADOBE_CLIENT_ID"],
            client_secret=os.environ["ADOBE_CLIENT_SECRET"],
            organization_id=os.environ["ADOBE_ORGANIZATION_ID"],
            account_id=os.environ["ADOBE_ACCOUNT_ID"],
        )
        pdf_services = PDFServices(credentials=creds)

        with open(input_path, "rb") as f:
            input_bytes = f.read()

        print(">>> [DEBUG] upload - 파일 업로드 시작", flush=True)
        try:
            asset = pdf_services.upload(input_bytes, PDFServicesMediaType.PDF)
            print(">>> [DEBUG] upload - 파일 업로드 성공", flush=True)
        except Exception as upload_error:
            print(f">>> [DEBUG] upload - 파일 업로드 실패: {upload_error}", flush=True)
            raise

        params = ExportPDFParams(ExportPDFTargetFormat.DOCX)
        job = ExportPDFJob(asset, params)

        print(">>> [DEBUG] submit - 작업 제출 시작", flush=True)
        try:
            location = pdf_services.submit(job)
            print(f">>> [DEBUG] submit - 작업 제출 성공, location: {location}", flush=True)
        except Exception as submit_error:
            print(f">>> [DEBUG] submit - 작업 제출 실패: {submit_error}", flush=True)
            raise

        print(">>> [DEBUG] get_job_result - 결과 대기 시작", flush=True)
        try:
            result_asset = pdf_services.get_job_result(location, ExportPDFResult)
            print(">>> [DEBUG] get_job_result - 결과 대기 성공", flush=True)
        except Exception as result_error:
            print(f">>> [DEBUG] get_job_result - 결과 대기 실패: {result_error}", flush=True)
            raise

        print(">>> [DEBUG] get_content - 콘텐츠 다운로드 시작", flush=True)
        try:
            content = pdf_services.get_content(result_asset)
            print(f">>> [DEBUG] get_content - 콘텐츠 다운로드 성공, 크기: {len(content)} bytes", flush=True)
        except Exception as content_error:
            print(f">>> [DEBUG] get_content - 콘텐츠 다운로드 실패: {content_error}", flush=True)
            raise

        with open(output_path, "wb") as f:
            f.write(content)

        print(">>> [DEBUG] success - 전체 변환 성공", flush=True)
        return True, info

    except Exception as e:
        info = {
            "type": type(e).__name__,
            "status": getattr(e, "status_code", None),
            "error_code": getattr(e, "error_code", None),
            "message": getattr(e, "message", str(e)),
            "request_id": getattr(e, "request_id", None),
            "error_report": getattr(e, "error_report", None),
        }
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", flush=True)
        print("❌ Adobe call failed", flush=True)
        print(f"파일: {input_path}", flush=True)
        print(f"파일 크기: {os.path.getsize(input_path) if os.path.exists(input_path) else 'N/A'} bytes", flush=True)
        for k, v in info.items():
            print(f"{k}: {v}", flush=True)
        
        # 400 에러에 대한 추가 분석
        if info.get("status") == 400:
            print(">>> [DEBUG] HTTP 400 에러 분석:", flush=True)
            print("  - 가능한 원인: 손상된 PDF, 암호화된 PDF, 지원되지 않는 PDF 형식", flush=True)
            print("  - 또는 Adobe API 요청 형식 오류", flush=True)
        
        traceback.print_exc()
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", flush=True)
        return False, info

def extract_text_with_layout_from_pdf(pdf_path: str) -> Dict[str, Any]:
    """PDF에서 레이아웃 정보와 함께 텍스트 추출 (pdfplumber 사용)"""
    try:
        import pdfplumber
        all_text_blocks = []
        
        # PDF 방향 분석
        orientation_info = analyze_pdf_orientation(pdf_path)
        print(f"PDF 방향 분석 결과: {orientation_info['primary_orientation']} (가로: {orientation_info['landscape_pages']}, 세로: {orientation_info['portrait_pages']})")
        
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # 텍스트 추출
                text = page.extract_text()
                if text:
                    lines = text.split('\n')
                    for line_num, line in enumerate(lines):
                        if line.strip():
                            # 간단한 정렬 감지 (왼쪽 정렬로 기본 설정)
                            alignment = 'left'
                            if line.strip().center(len(line)) == line:
                                alignment = 'center'
                            elif line.startswith(' ' * 10):  # 많은 공백으로 시작하면 오른쪽 정렬로 추정
                                alignment = 'right'
                            
                            all_text_blocks.append({
                                'text': clean_special_characters(line.strip()),
                                'bbox': [0, line_num * 12, page.width, (line_num + 1) * 12],  # 추정 bbox
                                'page': page_num,
                                'alignment': alignment
                            })
        
        return {
            'text_blocks': all_text_blocks,
            'full_text': '\n'.join([block['text'] for block in all_text_blocks]),
            'orientation_info': orientation_info
        }
        
    except Exception as e:
        print(f"PDF 레이아웃 추출 중 오류: {e}")
        return {'text_blocks': [], 'full_text': ''}

def extract_text_blocks_with_ocr(image):
    """OCR을 사용하여 이미지에서 텍스트 블록 추출 (개선된 버전)"""
    try:
        # OCR 가용성 확인
        try:
            import pytesseract
            # Tesseract 경로 자동 감지 (Render 환경 대응)
            if os.path.exists('/usr/bin/tesseract'):
                pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
        except ImportError:
            print("  - pytesseract를 사용할 수 없습니다.")
            return ""
        
        # 이미지 전처리로 OCR 정확도 향상
        img_array = np.array(image)
        
        # 그레이스케일 변환
        if len(img_array.shape) == 3:
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        else:
            gray = img_array
        
        # 노이즈 제거
        denoised = cv2.medianBlur(gray, 3)
        
        # 대비 향상
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(denoised)
        
        # OCR 수행
        try:
            config = r"--oem 3 --psm 6 -l kor+eng"
            text = pytesseract.image_to_string(enhanced, config=config)
        except Exception as ocr_error:
            print(f"  - 한국어+영어 OCR 실패: {ocr_error}")
            # Fallback: 영어만으로 재시도
            try:
                config = r"--oem 3 --psm 6 -l eng"
                text = pytesseract.image_to_string(enhanced, config=config)
                print("  - 영어 OCR로 fallback 성공")
            except Exception:
                print("  - OCR 완전 실패")
                return ""
        
        if text.strip():
            cleaned_text = clean_special_characters(text.strip())
            print(f"  - OCR 텍스트 추출됨: {len(cleaned_text)}자")
            return cleaned_text
        else:
            print("  - OCR에서 텍스트를 찾을 수 없음")
            return ""
            
    except Exception as e:
        print(f"  - OCR 처리 중 오류: {e}")
        return ""

def convert_pdf_to_docx_with_adobe_direct(pdf_path, output_path):
    """Adobe PDF Services API를 사용하여 PDF를 DOCX로 직접 변환하는 함수"""
    if not ADOBE_SDK_AVAILABLE:
        print("Adobe PDF Services SDK를 사용할 수 없습니다.")
        return False
    
    try:
        # Adobe API 자격증명 설정
        client_id = ADOBE_CONFIG["client_credentials"]["client_id"]
        client_secret = ADOBE_CONFIG["client_credentials"]["client_secret"]
        organization_id = ADOBE_CONFIG["service_principal_credentials"]["organization_id"]
        account_id = ADOBE_CONFIG["service_principal_credentials"]["account_id"]
        
        if not all([client_id, client_secret, organization_id, account_id]):
            print("Adobe API 자격증명이 완전하지 않습니다.")
            return False
        
        # ServicePrincipalCredentials 생성 (OAuth Server-to-Server)
        credentials = ServicePrincipalCredentials(
            client_id=client_id,
            client_secret=client_secret
        )
        
        # PDFServices 인스턴스 생성
        pdf_services = PDFServices(credentials=credentials)
        
        # 입력 파일을 StreamAsset으로 생성
        with open(pdf_path, 'rb') as file:
            input_stream = file.read()
        
        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PDF)
        
        # ExportPDF 작업 매개변수 설정
        export_pdf_params = ExportPDFParams(
            target_format=ExportPDFTargetFormat.DOCX
        )
        
        # ExportPDF 작업 생성
        export_pdf_job = ExportPDFJob(input_asset=input_asset, export_pdf_params=export_pdf_params)
        
        print(">>> [DEBUG 1] Adobe 변환 함수 진입")
        try:
            print(">>> [DEBUG 2] try 블록 진입, execute() 호출 직전")
            
            # 작업 제출 및 결과 대기 - 실제 Adobe API 실행 지점
            location = pdf_services.submit(export_pdf_job)
            pdf_services_response = pdf_services.get_job_result(location, ExportPDFResult)
            
            print(">>> [DEBUG 3] execute() 호출 성공")
            conversion_success = True  # 성공했음을 표시
            
        except ServiceApiException as e:
            # Adobe API 관련 에러 (가장 흔함)
            print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
            print(f"❌ Adobe ServiceApiException 발생: {e}")
            print(f"    - Status Code: {e.status_code if hasattr(e, 'status_code') else 'N/A'}")
            print(f"    - Error Code: {e.error_code if hasattr(e, 'error_code') else 'N/A'}")
            print(f"    - Error Message: {e.message if hasattr(e, 'message') else str(e)}")
            print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
            conversion_success = False
            raise  # 기존 예외 처리로 전달
            
        except Exception as e:
            # 그 외 모든 예상치 못한 에러
            print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
            print(f"❌ 변환 중 알 수 없는 예외 발생: {str(e)}")
            import traceback
            traceback.print_exc()
            print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
            conversion_success = False
            raise  # 기존 예외 처리로 전달
            
        print(">>> [DEBUG 4] Adobe 변환 함수 종료")
        
        # 결과 다운로드
        result_asset = pdf_services_response.get_result().get_asset()
        stream_asset = pdf_services.get_content(result_asset)
        
        # 결과를 파일로 저장
        with open(output_path, "wb") as file:
            file.write(stream_asset.get_input_stream())
        
        print(f"Adobe API를 사용하여 PDF를 DOCX로 성공적으로 변환했습니다: {output_path}")
        return True
        
    except ServiceApiException as e:
        print(f"Adobe ServiceApiException: {e}")
        print(f"Request ID: {getattr(e, 'request_id', 'N/A')}")
        print(f"Status Code: {getattr(e, 'status_code', 'N/A')}")
        print(f"Error Code: {getattr(e, 'error_code', 'N/A')}")
        print(f"Error Message: {getattr(e, 'message', str(e))}")
        return False
        
    except Exception as e:
        print(f"Adobe API 변환 중 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        return False

def adobe_pdf_to_docx(input_path, output_path):
    """사용자 요청에 따른 Adobe API PDF to DOCX 변환 함수"""
    print(">>> [DEBUG] adobe_pdf_to_docx 함수가 시작되었습니다.")  # <--- 이 줄을 함수 가장 처음에 추가!
    try:
        print(">>> [DEBUG] adobe_pdf_to_docx: try 블록 진입. SDK 초기화 시도.")
        
        # Adobe API 변환 실행
        result = convert_pdf_to_docx_with_adobe_direct(input_path, output_path)
        
        if result:
            print(">>> [DEBUG] adobe_pdf_to_docx: execute() 성공.")
            return True
        else:
            print(">>> [DEBUG] adobe_pdf_to_docx: execute() 실패.")
            return False
            
    except ServiceApiException as e:
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        print(f"❌ Adobe ServiceApiException: {e.message if hasattr(e, 'message') else str(e)}")
        print(f"    - Status Code: {e.status_code if hasattr(e, 'status_code') else 'N/A'}, Error Code: {e.error_code if hasattr(e, 'error_code') else 'N/A'}")
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        return False
        
    except Exception as e:
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        print(f"❌ adobe_pdf_to_docx에서 예외 발생: {str(e)}")
        import traceback
        traceback.print_exc()
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        return False

def extract_pdf_content_with_adobe(pdf_path):
    """Adobe PDF Services API를 사용하여 PDF 내용을 추출하는 함수"""
    if not ADOBE_SDK_AVAILABLE:
        print("Adobe PDF Services SDK를 사용할 수 없습니다.")
        return None
    
    # Adobe API 환경변수 확인
    client_id = ADOBE_CONFIG["client_credentials"]["client_id"]
    client_secret = ADOBE_CONFIG["client_credentials"]["client_secret"]
    
    if not client_id or not client_secret:
        print("Adobe API 환경변수가 설정되지 않았습니다. fallback 모드로 진행합니다.")
        return None
        
    try:
        # Adobe API 자격 증명 설정 (올바른 클래스 사용)
        credentials = ServicePrincipalCredentials(
            client_id=client_id,
            client_secret=client_secret
        )
        
        # PDF Services 인스턴스 생성
        pdf_services = PDFServices(credentials=credentials)
        
        # PDF 파일을 스트림으로 읽기
        with open(pdf_path, 'rb') as file:
            input_stream = file.read()
        
        # StreamAsset 생성
        input_asset = pdf_services.upload(input_stream=input_stream, mime_type=PDFServicesMediaType.PDF)
        
        print("Adobe API를 사용하여 PDF 내용을 처리했습니다.")
        return input_asset
            
    except ServiceApiException as e:
        # Adobe API 관련 에러 (가장 흔함)
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        print(f"❌ Adobe ServiceApiException 발생: {e}")
        print(f"    - Request ID: {getattr(e, 'request_id', 'N/A')}")
        print(f"    - Status Code: {getattr(e, 'status_code', 'N/A')}")
        print(f"    - Error Code: {getattr(e, 'error_code', 'N/A')}")
        print(f"    - Error Message: {getattr(e, 'message', str(e))}")
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        return None
    except (ServiceUsageException, SdkException) as e:
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        print(f"❌ Adobe SDK 오류 발생: {type(e).__name__}: {str(e)}")
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        return None
    except Exception as e:
        # 그 외 모든 예상치 못한 에러
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        print(f"❌ 변환 중 알 수 없는 예외 발생: {str(e)}")
        import traceback
        traceback.print_exc()
        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        return None

def pdf_to_docx(pdf_path, output_path, quality='medium'):
    """PDF를 DOCX로 변환하는 함수 (Adobe API 우선, pdf2docx 및 OCR 보조)"""
    try:
        # 파일명에서 확장자 제거하여 디버깅용 prefix 생성
        filename_prefix = os.path.splitext(os.path.basename(pdf_path))[0]
        
        # 1단계: Adobe API를 최우선으로 시도
        if adobe_available and is_adobe_api_available():
            print(">>> [DEBUG] Adobe 분기 진입. adobe_pdf_to_docx 함수 호출 시도.")
            try:
                conversion_success = adobe_pdf_to_docx(pdf_path, output_path)
                print(">>> [DEBUG] adobe_pdf_to_docx 함수 정상 종료. 결과:", conversion_success)
                if conversion_success:
                    return True
                else:
                    print("Adobe API 직접 변환 실패, Extract API로 시도...")
                    
            except Exception as e:
                print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
                print(f"❌ /convert 라우트에서 예외 발생: {str(e)}")
                import traceback
                traceback.print_exc()
                print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
                conversion_success = False
        else:
            print("⚠️ Adobe API 사용 불가 - fallback 모드로 변환합니다.")
            if not ADOBE_SDK_AVAILABLE:
                print("  - Adobe SDK가 설치되지 않음")
            else:
                print("  - Adobe API 환경변수가 설정되지 않음")
        
        # 2단계: pdf2docx 라이브러리를 두 번째로 시도
        print("=== 2단계: pdf2docx 라이브러리 변환 시도 ===")
        if pdf_to_docx_with_pdf2docx(pdf_path, output_path):
            print("pdf2docx 변환 성공! Microsoft Word 호환성 확인...")
            
            # 변환된 파일이 실제로 존재하고 크기가 적절한지 확인
            if os.path.exists(output_path) and os.path.getsize(output_path) > 1024:  # 1KB 이상
                print(f"변환 완료: {output_path} (크기: {os.path.getsize(output_path)} bytes)")
                return True
            else:
                print("pdf2docx 변환 결과가 부적절함. 대체 방법 시도...")
        
        print("=== 3단계: 기존 OCR 방법으로 fallback ===")
        # 품질 설정에 따른 파라미터 설정 (최적화됨)
        quality_settings = {
            'medium': {
                'dpi': 120,  # DPI 최적화로 속도 향상
                'format': 'jpeg',
                'jpeg_quality': 80,  # 품질과 속도의 균형
                'max_size': (1600, 1200),  # 적절한 해상도
                'description': '균형 변환 (최적화된 속도와 품질)'
            },
            'high': {
                'dpi': 180,  # 고품질이지만 속도 고려
                'format': 'jpeg',  # PNG 대신 JPEG 사용으로 속도 향상
                'jpeg_quality': 90,
                'max_size': (2048, 1536),  # 해상도 최적화
                'description': '고품질 변환 (향상된 속도)'
            }
        }
        
        settings = quality_settings.get(quality, quality_settings['medium'])
        print(f"변환 설정: {settings['description']}")
        
        # 1단계: 레이아웃 인식을 통한 텍스트 추출 시도
        print("레이아웃 인식을 통한 텍스트 추출을 시도합니다...")
        layout_data = extract_text_with_layout_from_pdf(pdf_path)
        extracted_text = layout_data.get('full_text', '')
        text_blocks = layout_data.get('text_blocks', [])
        orientation_info = layout_data.get('orientation_info', {})
        
        if extracted_text:
            print(f"레이아웃 인식으로 텍스트 추출 성공: {len(extracted_text)}자")
        else:
            print("레이아웃 인식 실패, Adobe API를 시도합니다...")
            
            # 2단계: Adobe API를 사용한 PDF 내용 추출 시도
            if adobe_available and is_adobe_api_available():
                print("✅ Adobe API로 변환을 시작합니다.")
                try:
                    # 이 부분이 실제 Adobe SDK를 사용하는 코드입니다.
                    extracted_content = extract_pdf_content_with_adobe(pdf_path)
                    if extracted_content:
                        extracted_text = str(extracted_content)
                        print(f"Adobe API에서 텍스트 추출 성공: {len(extracted_text)}자")
                    else:
                        print("Adobe API 추출 실패, OCR 방법으로 진행합니다.")
                        
                except ServiceApiException as e:
                    # Adobe API 관련 에러 (가장 흔함)
                    print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
                    print(f"❌ Adobe ServiceApiException 발생: {e}")
                    print(f"    - Request ID: {getattr(e, 'request_id', 'N/A')}")
                    print(f"    - Status Code: {getattr(e, 'status_code', 'N/A')}")
                    print(f"    - Error Code: {getattr(e, 'error_code', 'N/A')}")
                    print(f"    - Error Message: {getattr(e, 'message', str(e))}")
                    print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
                    extracted_text = ''
                    
                except Exception as e:
                    # 그 외 모든 예상치 못한 에러
                    print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
                    print(f"❌ 변환 중 알 수 없는 예외 발생: {str(e)}")
                    import traceback
                    traceback.print_exc()
                    print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
                    extracted_text = ''
            else:
                # Adobe API를 사용할 수 없을 때의 대체 로직
                print("⚠️ Adobe API 사용 불가 - fallback 모드로 변환합니다.")
        
        # 기본 방법: PDF를 이미지로 변환 (품질별 최적화)
        print("PDF를 이미지로 변환 중...")
        images = convert_from_path(pdf_path, dpi=settings['dpi'], fmt=settings['format'])
        
        # 디버깅: 변환된 이미지들을 저장
        print("=== 디버깅: 변환된 이미지 저장 ===")
        for i, image in enumerate(images):
            save_debug_image(image, filename_prefix, i+1)
        
        # 새 Word 문서 생성 - 호환성 개선 및 방향 자동 감지
        doc = Document()
        
        # 페이지 설정 (문서 방향에 따라 자동 조정)
        section = doc.sections[0]
        primary_orientation = orientation_info.get('primary_orientation', 'portrait')
        
        if primary_orientation == 'landscape':
            # 가로형 문서 설정
            section.page_width = Inches(11)
            section.page_height = Inches(8.5)
            section.left_margin = Inches(0.8)
            section.right_margin = Inches(0.8)
            section.top_margin = Inches(0.6)
            section.bottom_margin = Inches(0.6)
            print("가로형 문서로 설정됨")
        else:
            # 세로형 문서 설정
            section.page_width = Inches(8.5)
            section.page_height = Inches(11)
            section.left_margin = Inches(1)
            section.right_margin = Inches(1)
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            print("세로형 문서로 설정됨")
        
        # 문서 속성 설정 (Microsoft Word 호환성 향상)
        try:
            # 안전한 파일명 생성 (특수문자 제거)
            safe_filename = re.sub(r'[^\w\s-]', '', os.path.splitext(os.path.basename(pdf_path))[0])
            doc.core_properties.title = safe_filename[:50]  # 제목 길이 제한
            doc.core_properties.author = "Document Converter"
            doc.core_properties.subject = "PDF to DOCX Conversion"
            doc.core_properties.comments = "Converted using advanced OCR and layout recognition"
        except Exception as e:
            print(f"문서 속성 설정 중 오류 (무시됨): {e}")
        
        all_ocr_text = []
        
        print(f"총 {len(images)}페이지 처리 중...")
        # OCR로 텍스트 추출 (Adobe API가 실패한 경우)
        if not extracted_text:
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} OCR 처리 중...")
                ocr_text = extract_text_blocks_with_ocr(image)
                if ocr_text:
                    all_ocr_text.append(ocr_text)
        
        # 편집 가능한 텍스트만 추가 (원본 이미지 제거)
        final_text = extracted_text if extracted_text else '\n'.join(all_ocr_text)
        
        # 디버깅: 추출된 텍스트를 파일로 저장
        print("=== 디버깅: 추출된 텍스트 저장 ===")
        if final_text:
            save_debug_text(final_text, filename_prefix)
        elif all_ocr_text:
            combined_ocr_text = '\n'.join(all_ocr_text)
            save_debug_text(combined_ocr_text, filename_prefix + "_ocr")
        else:
            save_debug_text("텍스트 추출 실패", filename_prefix + "_no_text")
        
        # --- 핵심 수정 부분: OCR 텍스트 추출 실패 시 None 반환 ---
        if not final_text.strip() and not text_blocks:
            print(f"'{pdf_path}' 파일에서 유효한 텍스트를 찾지 못했습니다.")
            print(f"OCR 텍스트 길이: {len(final_text)}, 텍스트 블록 수: {len(text_blocks) if text_blocks else 0}")
            print(f"이미지 품질이 낮거나 텍스트가 없는 PDF 파일일 수 있습니다: {pdf_path}")
            print(f"변환 품질 설정: {quality}, 이미지 수: {len(images) if 'images' in locals() else 'N/A'}")
            return None  # 텍스트가 없으면 None을 반환
        
        if final_text or text_blocks:
            print(f"편집 가능한 텍스트 문서 생성: {len(final_text)}자")
            
            # 레이아웃 정보가 있는 경우 페이지별로 구성 (페이지 헤더 제거)
            if text_blocks:
                print("레이아웃 정보를 활용하여 텍스트 구조화...")
                
                # 페이지별로 텍스트 구성 (페이지 번호 헤더 없이)
                for page_num in range(len(images)):
                    if page_num > 0:
                        doc.add_page_break()
                    
                    # 해당 페이지의 텍스트 블록 추가
                    page_text_blocks = [block for block in text_blocks if block['page'] == page_num]
                    
                    if page_text_blocks:
                        for block in page_text_blocks:
                            text_paragraph = doc.add_paragraph()
                            text_run = text_paragraph.add_run(block['text'])
                            
                            # 텍스트 정렬 적용 (원본과 동일하게)
                            if block['alignment'] == 'center':
                                text_paragraph.alignment = 1  # 중앙 정렬
                                text_run.bold = True  # 중앙 정렬 텍스트는 굵게
                            elif block['alignment'] == 'right':
                                text_paragraph.alignment = 2  # 오른쪽 정렬
                            else:
                                text_paragraph.alignment = 0  # 왼쪽 정렬
                            
                            # 원본과 동일한 텍스트 스타일 적용 (Microsoft Word 호환성 개선)
                            try:
                                # 폰트 설정 (한글 문서에 적합한 폰트)
                                text_run.font.name = '맑은 고딕'
                                text_run.font.size = Pt(11)  # 표준 문서 크기
                                
                                # 제목 스타일 적용
                                if len(block['text']) < 50 and block['alignment'] == 'center':
                                    text_run.font.size = Pt(14)
                                    text_run.bold = True
                                elif '제목' in block['text'] or '공문' in block['text']:
                                    text_run.font.size = Pt(13)
                                    text_run.bold = True
                                
                                # 줄간격 및 단락 간격 설정 (원본과 동일하게)
                                text_paragraph.paragraph_format.line_spacing = 1.2
                                text_paragraph.paragraph_format.space_after = Pt(3)
                                text_paragraph.paragraph_format.space_before = Pt(0)
                                
                                # 들여쓰기 설정 (원본 레이아웃 유지)
                                if block['alignment'] == 'left':
                                    text_paragraph.paragraph_format.left_indent = Pt(0)
                                elif block['alignment'] == 'center':
                                    text_paragraph.paragraph_format.left_indent = Pt(0)
                                    
                            except Exception as e:
                                print(f"텍스트 스타일 설정 중 오류 (무시됨): {e}")
            else:
                # 레이아웃 정보가 없는 경우 일반 텍스트로 추가 (Microsoft Word 호환성 개선)
                clean_final_text = final_text.replace('\x00', '').replace('\ufffd', '').strip()
                if clean_final_text:
                    paragraphs = clean_final_text.split('\n\n')
                    for para_text in paragraphs:
                        if para_text.strip():
                            text_paragraph = doc.add_paragraph()
                            text_run = text_paragraph.add_run(para_text.strip())
                            
                            # 일반 텍스트에도 표준 스타일 적용
                            try:
                                text_run.font.name = '맑은 고딕'
                                text_run.font.size = Pt(11)
                                text_paragraph.paragraph_format.line_spacing = 1.2
                                text_paragraph.paragraph_format.space_after = Pt(3)
                                text_paragraph.paragraph_format.space_before = Pt(0)
                            except Exception as e:
                                print(f"일반 텍스트 스타일 설정 중 오류 (무시됨): {e}")
                else:
                    text_paragraph = doc.add_paragraph()
                    text_run = text_paragraph.add_run("텍스트를 추출할 수 없었습니다.")
                    try:
                        text_run.font.name = '맑은 고딕'
                        text_run.font.size = Pt(11)
                    except Exception as e:
                        print(f"오류 메시지 스타일 설정 중 오류 (무시됨): {e}")
            
            print("편집 가능한 텍스트 문서가 생성되었습니다.")
        else:
            print("추출할 수 있는 텍스트가 없습니다. 이미지 기반 문서를 생성합니다.")
            
            # 텍스트가 없는 경우에만 이미지 추가
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} 처리 중...")
                
                # 이미지 크기 최적화 (원본 문서와 동일한 크기 유지)
                original_width, original_height = image.size
                
                # 문서 방향에 따른 이미지 크기 조정
                if primary_orientation == 'landscape':
                    # 가로형 문서: 최대 너비 10인치
                    target_width = min(10, original_width / 72)  # 72 DPI 기준
                    aspect_ratio = original_height / original_width
                    target_height = target_width * aspect_ratio
                else:
                    # 세로형 문서: 최대 너비 6.5인치
                    target_width = min(6.5, original_width / 72)  # 72 DPI 기준
                    aspect_ratio = original_height / original_width
                    target_height = target_width * aspect_ratio
                
                # 이미지를 임시 파일로 저장 (JPEG 최적화)
                temp_img_path = None
                try:
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                        temp_img_path = temp_img.name
                        # JPEG로 저장 (품질별 압축, 빠른 처리)
                        image.save(temp_img_path, 'JPEG', quality=settings['jpeg_quality'], optimize=True)
                    
                    # 문서에 이미지 추가 (원본 비율 유지)
                    doc.add_picture(temp_img_path, width=DocxInches(target_width))
                    
                    # 페이지 구분을 위한 페이지 브레이크 추가 (마지막 페이지 제외)
                    if i < len(images) - 1:
                        doc.add_page_break()
                    
                finally:
                    # 임시 파일 삭제 (빠른 처리)
                    if temp_img_path and os.path.exists(temp_img_path):
                        try:
                            os.unlink(temp_img_path)
                        except (OSError, PermissionError) as e:
                            print(f"임시 파일 삭제 실패 (무시됨): {e}")
        
        # DOCX 파일 저장 (Microsoft Word 호환성 최적화)
        try:
            # 임시 파일로 먼저 저장 후 이동 (안전한 저장)
            temp_output = output_path + '.tmp'
            doc.save(temp_output)
            
            # 기존 파일이 있으면 삭제
            if os.path.exists(output_path):
                os.remove(output_path)
            
            # 임시 파일을 최종 파일로 이동
            os.rename(temp_output, output_path)
            
            print(f"DOCX 파일 저장 완료: {output_path}")
            print("Microsoft Word 호환성이 개선된 문서가 생성되었습니다.")
            return True
            
        except Exception as save_error:
            print(f"DOCX 파일 저장 중 오류: {save_error}")
            # 임시 파일 정리
            temp_output = output_path + '.tmp'
            if os.path.exists(temp_output):
                try:
                    os.remove(temp_output)
                except:
                    pass
            return False
        
    except Exception as e:
        print(f"변환 중 오류 발생: {str(e)}")
        return False

def pdf_to_pptx(pdf_path, output_path, quality='medium'):
    """PDF를 PPTX로 변환하는 함수 (Adobe API 통합 및 OCR 텍스트 추출, 방향 자동 감지)"""
    try:
        # 품질 설정에 따른 파라미터 설정 (최적화됨)
        quality_settings = {
            'medium': {
                'dpi': 120,  # DPI 최적화로 속도 향상
                'format': 'jpeg',
                'jpeg_quality': 80,  # 품질과 속도의 균형
                'max_size': (1600, 1200),  # 적절한 해상도
                'description': '균형 변환 (최적화된 속도와 품질)'
            },
            'high': {
                'dpi': 180,  # 고품질이지만 속도 고려
                'format': 'jpeg',  # PNG 대신 JPEG 사용으로 속도 향상
                'jpeg_quality': 90,
                'max_size': (2048, 1536),  # 해상도 최적화
                'description': '고품질 변환 (향상된 속도)'
            }
        }
        
        settings = quality_settings.get(quality, quality_settings['medium'])
        print(f"변환 설정: {settings['description']}")
        
        # 1단계: 레이아웃 인식을 통한 텍스트 추출 시도 (방향 정보 포함)
        print("레이아웃 인식을 통한 텍스트 추출을 시도합니다...")
        layout_data = extract_text_with_layout_from_pdf(pdf_path)
        extracted_text = layout_data.get('full_text', '')
        text_blocks = layout_data.get('text_blocks', [])
        orientation_info = layout_data.get('orientation_info', {})
        all_ocr_text = []
        
        if extracted_text:
            print(f"레이아웃 인식으로 텍스트 추출 성공: {len(extracted_text)}자")
        else:
            print("레이아웃 인식 실패, Adobe API를 시도합니다...")
            
            # 2단계: Adobe API를 사용한 PDF 내용 추출 시도
            if adobe_available:
                print("Adobe API를 사용하여 PDF 처리를 시작합니다...")
                extracted_content = extract_pdf_content_with_adobe(pdf_path)
                if extracted_content:
                    extracted_text = str(extracted_content)
                    print(f"Adobe API에서 텍스트 추출 성공: {len(extracted_text)}자")
                else:
                    print("Adobe API 추출 실패, OCR 방법으로 진행합니다.")
        
        # 기본 방법: PDF를 이미지로 변환 (품질별 최적화)
        print("PDF를 이미지로 변환 중...")
        images = convert_from_path(pdf_path, dpi=settings['dpi'], fmt=settings['format'])
        
        # 새 PowerPoint 프레젠테이션 생성 (방향에 따른 슬라이드 설정)
        prs = Presentation()
        
        # 슬라이드 크기 설정 (문서 방향에 따라 자동 조정)
        primary_orientation = orientation_info.get('primary_orientation', 'portrait')
        
        if primary_orientation == 'landscape':
            # 가로형 슬라이드 (16:9 비율)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            print("가로형 슬라이드로 설정됨 (16:9)")
        else:
            # 세로형 슬라이드 (9:16 비율)
            prs.slide_width = Inches(7.5)
            prs.slide_height = Inches(13.33)
            print("세로형 슬라이드로 설정됨 (9:16)")
        
        all_ocr_text = []
        
        print(f"총 {len(images)}페이지 처리 중...")
        def get_blank_slide_layout(prs):
            """안전한 빈 슬라이드 레이아웃 가져오기"""
            try:
                # 슬라이드 레이아웃이 있는지 먼저 확인
                if len(prs.slide_layouts) == 0:
                    raise IndexError("슬라이드 레이아웃이 없습니다")
                
                # 빈 슬라이드 레이아웃 우선 선택
                if len(prs.slide_layouts) > 6:
                    return prs.slide_layouts[6]  # 빈 슬라이드
                elif len(prs.slide_layouts) > 5:
                    return prs.slide_layouts[5]  # 제목만 있는 슬라이드
                else:
                    return prs.slide_layouts[0]  # 첫 번째 사용 가능한 레이아웃
            except (IndexError, AttributeError) as e:
                print(f"슬라이드 레이아웃 접근 오류: {e}")
                # 기본 프레젠테이션 생성 시 최소 하나의 레이아웃은 있어야 함
                if len(prs.slide_layouts) > 0:
                    return prs.slide_layouts[0]
                else:
                    raise Exception("사용 가능한 슬라이드 레이아웃이 없습니다")
        
        # 편집 가능한 텍스트 슬라이드 생성 (원본 이미지 제거)
        # OCR로 텍스트 추출 (Adobe API가 실패한 경우)
        if not extracted_text:
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} OCR 처리 중...")
                ocr_text = extract_text_blocks_with_ocr(image)
                if ocr_text:
                    all_ocr_text.append(ocr_text)
        
        # 편집 가능한 텍스트 슬라이드 생성
        final_text = extracted_text if extracted_text else '\n'.join(all_ocr_text)
        
        if text_blocks:
            print(f"편집 가능한 텍스트 슬라이드 생성: {len(text_blocks)}개 블록")
            
            # 페이지별로 슬라이드 구성
            for page_num in range(len(images)):
                # 새 슬라이드 추가 (제목과 내용 레이아웃)
                try:
                    slide_layout = prs.slide_layouts[1]  # 제목과 내용 레이아웃
                except (IndexError, AttributeError):
                    slide_layout = get_blank_slide_layout(prs)
                
                slide = prs.slides.add_slide(slide_layout)
                
                # 슬라이드 제목 설정
                try:
                    title_shape = slide.shapes.title
                    title_shape.text = f"페이지 {page_num + 1}"
                except AttributeError:
                    # 제목이 없는 레이아웃인 경우 텍스트 박스 추가
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(1)
                    title_box = slide.shapes.add_textbox(left, top, width, height)
                    title_frame = title_box.text_frame
                    title_frame.text = f"페이지 {page_num + 1}"
                
                # 해당 페이지의 텍스트 블록 추가
                page_text_blocks = [block for block in text_blocks if block['page'] == page_num]
                
                if page_text_blocks:
                    # 내용 텍스트박스 가져오기
                    try:
                        content_shape = slide.placeholders[1]
                        text_frame = content_shape.text_frame
                        text_frame.clear()
                        
                        for j, block in enumerate(page_text_blocks):
                            if j == 0:
                                # 첫 번째 단락
                                p = text_frame.paragraphs[0]
                            else:
                                # 추가 단락
                                p = text_frame.add_paragraph()
                            
                            p.text = block['text']
                            
                            # 텍스트 정렬 적용
                            if block['alignment'] == 'center':
                                p.alignment = 1  # 중앙 정렬
                                try:
                                    p.font.bold = True
                                except AttributeError:
                                    pass
                            elif block['alignment'] == 'right':
                                p.alignment = 2  # 오른쪽 정렬
                            else:
                                p.alignment = 0  # 왼쪽 정렬
                            
                            # 텍스트 크기 조정
                            try:
                                from docx.shared import Pt
                                if len(block['text']) < 50 and block['alignment'] == 'center':
                                    p.font.size = Pt(18)  # 제목용 크기
                                else:
                                    p.font.size = Pt(14)  # 본문용 크기
                            except (ImportError, AttributeError):
                                pass
                                
                    except (IndexError, AttributeError):
                        # 내용 플레이스홀더가 없는 경우 텍스트 박스 추가
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(9)
                        height = Inches(6)
                        content_box = slide.shapes.add_textbox(left, top, width, height)
                        content_frame = content_box.text_frame
                        content_text = '\n'.join([block['text'] for block in page_text_blocks])
                        content_frame.text = content_text
                else:
                    # 텍스트가 없는 페이지
                    try:
                        content_shape = slide.placeholders[1]
                        content_shape.text = "[이 페이지는 텍스트 추출이 어려운 이미지 페이지입니다]"
                    except (IndexError, AttributeError):
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(9)
                        height = Inches(6)
                        content_box = slide.shapes.add_textbox(left, top, width, height)
                        content_frame = content_box.text_frame
                        content_frame.text = "[이 페이지는 텍스트 추출이 어려운 이미지 페이지입니다]"
        
        elif final_text:
            print(f"일반 텍스트 슬라이드 생성: {len(final_text)}자")
            
            # 텍스트를 적절한 크기로 나누어 슬라이드 생성
            text_chunks = final_text.split('\n\n')
            chunk_size = 5  # 슬라이드당 단락 수
            
            for i in range(0, len(text_chunks), chunk_size):
                try:
                    slide_layout = prs.slide_layouts[1]  # 제목과 내용 레이아웃
                except (IndexError, AttributeError):
                    slide_layout = get_blank_slide_layout(prs)
                
                slide = prs.slides.add_slide(slide_layout)
                
                # 슬라이드 제목
                try:
                    title_shape = slide.shapes.title
                    title_shape.text = f"슬라이드 {(i // chunk_size) + 1}"
                except AttributeError:
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(1)
                    title_box = slide.shapes.add_textbox(left, top, width, height)
                    title_frame = title_box.text_frame
                    title_frame.text = f"슬라이드 {(i // chunk_size) + 1}"
                
                # 내용 추가
                chunk_text = text_chunks[i:i+chunk_size]
                content_text = '\n\n'.join([para.strip() for para in chunk_text if para.strip()])
                
                try:
                    content_shape = slide.placeholders[1]
                    content_shape.text = content_text
                except (IndexError, AttributeError):
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(6)
                    content_box = slide.shapes.add_textbox(left, top, width, height)
                    content_frame = content_box.text_frame
                    content_frame.text = content_text
        
        else:
            print("추출할 수 있는 텍스트가 없습니다. 이미지 기반 슬라이드를 생성합니다.")
            
            # 텍스트가 없는 경우에만 이미지 슬라이드 생성
            for i, image in enumerate(images):
                print(f"페이지 {i+1}/{len(images)} 처리 중...")
                
                # 슬라이드 추가 - 안전한 레이아웃 사용
                slide_layout = get_blank_slide_layout(prs)
                slide = prs.slides.add_slide(slide_layout)
                
                # 이미지 크기 최적화 (원본 문서와 동일한 크기 유지)
                original_width, original_height = image.size
                
                # 슬라이드 방향에 따른 이미지 크기 조정
                if primary_orientation == 'landscape':
                    # 가로형 슬라이드: 최대 높이 6.5인치
                    target_height = min(6.5, original_height / 72)  # 72 DPI 기준
                    aspect_ratio = original_width / original_height
                    target_width = target_height * aspect_ratio
                    # 슬라이드 너비를 초과하지 않도록 조정
                    max_slide_width = 12.5  # 가로형 슬라이드 최대 너비
                    if target_width > max_slide_width:
                        target_width = max_slide_width
                        target_height = target_width / aspect_ratio
                else:
                    # 세로형 슬라이드: 최대 너비 6.5인치
                    target_width = min(6.5, original_width / 72)  # 72 DPI 기준
                    aspect_ratio = original_height / original_width
                    target_height = target_width * aspect_ratio
                    # 슬라이드 높이를 초과하지 않도록 조정
                    max_slide_height = 12.5  # 세로형 슬라이드 최대 높이
                    if target_height > max_slide_height:
                        target_height = max_slide_height
                        target_width = target_height / aspect_ratio
                
                # 이미지를 임시 파일로 저장 (JPEG 최적화)
                temp_img_path = None
                try:
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
                        temp_img_path = temp_img.name
                        # JPEG로 저장 (품질별 압축, 빠른 처리)
                        image.save(temp_img_path, 'JPEG', quality=settings['jpeg_quality'], optimize=True)
                    
                    # 슬라이드에 이미지 추가 (원본 비율 유지, 중앙 배치)
                    left = Inches((13.33 - target_width) / 2) if primary_orientation == 'landscape' else Inches((7.5 - target_width) / 2)
                    top = Inches((7.5 - target_height) / 2) if primary_orientation == 'landscape' else Inches((13.33 - target_height) / 2)
                    slide.shapes.add_picture(temp_img_path, left, top, width=Inches(target_width), height=Inches(target_height))
                    
                finally:
                    # 임시 파일 삭제 (빠른 처리)
                    if temp_img_path and os.path.exists(temp_img_path):
                        try:
                            os.unlink(temp_img_path)
                        except (OSError, PermissionError) as e:
                            print(f"임시 파일 삭제 실패 (무시됨): {e}")
                            # 임시 파일 삭제 실패는 무시하고 계속 진행
        
        # 하이브리드 변환: 추출된 텍스트를 편집 가능한 형태로 마지막 슬라이드에 추가
        final_text = extracted_text if extracted_text else '\n'.join(all_ocr_text)
        if final_text:
            print(f"하이브리드 변환: 추출된 텍스트를 편집 가능한 형태로 추가: {len(final_text)}자")
            
            # 텍스트 전용 슬라이드 추가 - 안전한 레이아웃 선택
            try:
                # 제목과 내용 레이아웃이 있는지 확인
                if len(prs.slide_layouts) > 1:
                    text_slide_layout = prs.slide_layouts[1]  # 제목과 내용 레이아웃
                else:
                    text_slide_layout = get_blank_slide_layout(prs)
            except (IndexError, AttributeError):
                text_slide_layout = get_blank_slide_layout(prs)
            
            text_slide = prs.slides.add_slide(text_slide_layout)
            
            # 제목 설정 (안전한 방법)
            try:
                title = text_slide.shapes.title
                title.text = "추출된 텍스트 (편집 가능)"
            except AttributeError:
                # 제목이 없는 레이아웃인 경우 텍스트 박스 추가
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(1)
                title_box = text_slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_frame.text = "추출된 텍스트 (편집 가능)"
            
            # 내용 설정 (안전한 방법) - 레이아웃 정보 활용
            if text_blocks:
                print("레이아웃 정보를 활용하여 텍스트 구조화...")
                try:
                    content = text_slide.placeholders[1]
                    content_frame = content.text_frame
                    content_frame.clear()
                    
                    current_page = -1
                    for block in text_blocks:
                        # 페이지가 바뀌면 구분선 추가
                        if block['page'] != current_page:
                            if current_page != -1:
                                p = content_frame.add_paragraph()
                                p.text = f"\n--- 페이지 {block['page'] + 1} ---"
                            current_page = block['page']
                        
                        # 텍스트 단락 추가
                        p = content_frame.add_paragraph()
                        p.text = block['text']
                        
                        # 정렬 설정 (기본값으로 처리)
                        if block['alignment'] == 'center':
                            p.alignment = 1  # 중앙 정렬
                        elif block['alignment'] == 'right':
                            p.alignment = 2  # 오른쪽 정렬
                        else:
                            p.alignment = 0  # 왼쪽 정렬
                            
                except (IndexError, AttributeError):
                    # 내용 플레이스홀더가 없는 경우 텍스트 박스 추가
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(6)
                    content_box = text_slide.shapes.add_textbox(left, top, width, height)
                    content_frame = content_box.text_frame
                    content_frame.text = final_text
            else:
                # 레이아웃 정보가 없는 경우 일반 텍스트로 추가
                try:
                    content = text_slide.placeholders[1]
                    content.text = final_text
                except (IndexError, AttributeError):
                    # 내용 플레이스홀더가 없는 경우 텍스트 박스 추가
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(6)
                    content_box = text_slide.shapes.add_textbox(left, top, width, height)
                    content_frame = content_box.text_frame
                    content_frame.text = final_text
            
            print("편집 가능한 텍스트가 레이아웃 정보와 함께 슬라이드에 추가되었습니다.")
        else:
            print("추출할 수 있는 텍스트가 없습니다.")
        
        # PPTX 파일 저장
        prs.save(output_path)
        return True
        
    except Exception as e:
        print(f"변환 중 오류 발생: {str(e)}")
        return False

def image_to_docx(image_path, output_path):
    """이미지를 DOCX로 변환하는 함수"""
    try:
        print(f"이미지 → DOCX 변환 시작: {image_path} -> {output_path}")
        
        # 새 Word 문서 생성
        doc = Document()
        
        # 이미지 열기 및 크기 확인
        with Image.open(image_path) as img:
            # 이미지 크기 (픽셀)
            img_width, img_height = img.size
            print(f"원본 이미지 크기: {img_width} x {img_height} 픽셀")
            
            # A4 페이지 크기 (인치)
            page_width = 8.27  # A4 너비 (인치)
            page_height = 11.69  # A4 높이 (인치)
            margin = 1.0  # 여백 (인치)
            
            # 사용 가능한 영역
            available_width = page_width - (2 * margin)
            available_height = page_height - (2 * margin)
            
            # 이미지 비율 계산
            img_ratio = img_width / img_height
            available_ratio = available_width / available_height
            
            # 이미지 크기 조정 계산
            if img_ratio > available_ratio:
                # 이미지가 더 넓음 - 너비에 맞춤
                docx_width = DocxInches(available_width)
                docx_height = DocxInches(available_width / img_ratio)
            else:
                # 이미지가 더 높음 - 높이에 맞춤
                docx_height = DocxInches(available_height)
                docx_width = DocxInches(available_height * img_ratio)
            
            print(f"DOCX 이미지 크기: {docx_width.inches:.2f} x {docx_height.inches:.2f} 인치")
        
        # 문서에 이미지 추가
        paragraph = doc.add_paragraph()
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.add_picture(image_path, width=docx_width, height=docx_height)
        
        # 이미지를 중앙 정렬
        paragraph.alignment = 1  # 중앙 정렬
        
        # 문서 저장
        doc.save(output_path)
        print(f"이미지 → DOCX 변환 완료: {output_path}")
        
        return True
        
    except Exception as e:
        print(f"이미지 → DOCX 변환 중 오류: {str(e)}")
        return False

def docx_to_pdf(docx_path, output_path):
    """DOCX를 PDF로 변환하는 함수"""
    try:
        # Windows에서 LibreOffice 사용
        if platform.system() == "Windows":
            # LibreOffice 경로 찾기
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                "soffice"  # PATH에 있는 경우
            ]
            
            libreoffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path) or path == "soffice":
                    libreoffice_path = path
                    break
            
            if libreoffice_path:
                # LibreOffice를 사용하여 변환
                output_dir = os.path.dirname(output_path)
                cmd = [
                    libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", output_dir,
                    docx_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode == 0:
                    # 생성된 PDF 파일명 확인 및 이동
                    base_name = os.path.splitext(os.path.basename(docx_path))[0]
                    generated_pdf = os.path.join(output_dir, base_name + ".pdf")
                    
                    if os.path.exists(generated_pdf) and generated_pdf != output_path:
                        os.rename(generated_pdf, output_path)
                    
                    return os.path.exists(output_path)
                else:
                    print(f"LibreOffice 변환 실패: {result.stderr}")
                    return False
            else:
                print("LibreOffice를 찾을 수 없습니다.")
                return False
        else:
            print("현재 Linux/Mac에서의 DOCX → PDF 변환은 지원되지 않습니다.")
            return False
            
    except Exception as e:
        print(f"DOCX → PDF 변환 중 오류: {str(e)}")
        return False

# 파일 크기 초과 오류 처리
@app.errorhandler(413)
def too_large(e):
    flash('파일 크기가 100MB를 초과합니다. 더 작은 파일을 선택해주세요.')
    return redirect(url_for('index'))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/convert', methods=['POST'])
def convert_file_api():
    """API 방식의 파일 변환 엔드포인트"""
    try:
        if ENABLE_DEBUG_LOGS:
            print("파일 업로드 요청 시작")
            print(f"Request files: {request.files}")
            print(f"Request form: {request.form}")
            print(f"Request content type: {request.content_type}")
        
        # 환경변수 기반 설정 확인
        adobe_ready = is_adobe_api_available()
        conversion_method = "Adobe API" if adobe_ready else "pdf2docx + OCR"
        if ENABLE_DEBUG_LOGS:
            print(f"사용할 변환 방법: {conversion_method}")
        
        # 1단계: 파일 존재 여부 확인
        if 'file' not in request.files:
            return jsonify({
                'success': False, 
                'error': '파일이 선택되지 않았습니다.',
                'conversion_method': conversion_method
            }), 400
        
        file = request.files['file']
        
        # 2단계: 파일명 확인
        if not file or file.filename == '' or file.filename is None:
            return jsonify({
                'success': False, 
                'error': '파일이 선택되지 않았습니다.',
                'conversion_method': conversion_method
            }), 400
        
        # 3단계: 파일 내용 및 크기 확인
        try:
            file.seek(0, 2)
            file_size = file.tell()
            file.seek(0)
            
            if file_size == 0:
                return jsonify({
                    'success': False, 
                    'error': '업로드된 파일이 비어있습니다.',
                    'conversion_method': conversion_method
                }), 400
            
            if file_size < 10:
                return jsonify({
                    'success': False, 
                    'error': '파일이 너무 작습니다.',
                    'conversion_method': conversion_method
                }), 400
            
            max_size = MAX_FILE_SIZE_MB * 1024 * 1024
            if file_size > max_size:
                return jsonify({
                    'success': False, 
                    'error': f'파일 크기가 너무 큽니다. (최대: {max_size//(1024*1024)}MB)',
                    'conversion_method': conversion_method
                }), 400
            
            print(f"파일 크기: {file_size // (1024*1024) if file_size >= 1024*1024 else file_size // 1024}{'MB' if file_size >= 1024*1024 else 'KB'}")
            
            # 4단계: 파일 형식 검증
            file_content = file.read(10)
            file.seek(0)
            
            file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
            
            # 파일 형식별 검증
            if file_ext == 'pdf' and not file_content.startswith(b'%PDF-'):
                return jsonify({
                    'success': False, 
                    'error': '올바른 PDF 파일이 아닙니다.',
                    'conversion_method': conversion_method
                }), 400
            
        except Exception as e:
            print(f"파일 검증 중 오류: {str(e)}")
            return jsonify({
                'success': False, 
                'error': f'파일 검증 중 오류가 발생했습니다: {str(e)}',
                'conversion_method': conversion_method
            }), 400
        
        # 5단계: 파일 처리
        if file and allowed_file(file.filename):
            # 원본 파일명에서 확장자 추출
            original_filename = file.filename
            print(f"🔍 원본 파일명: '{original_filename}'")
            
            # 확장자 추출 (원본 파일명에서)
            if '.' not in original_filename:
                return jsonify({
                    'success': False, 
                    'error': '파일 확장자가 없습니다.',
                    'conversion_method': conversion_method
                }), 400
            
            file_ext = original_filename.rsplit('.', 1)[1].lower()
            print(f"🔍 추출된 확장자: '{file_ext}'")
            
            # 안전한 파일명 생성 (한글 파일명 처리)
            import re
            import unicodedata
            
            # 원본 파일명에서 확장자 제거
            base_name = original_filename.rsplit('.', 1)[0]
            
            # 한글 파일명을 안전하게 처리
            # 1. 유니코드 정규화
            normalized_name = unicodedata.normalize('NFC', base_name)
            # 2. 안전하지 않은 문자 제거 (경로 구분자, 특수문자 등)
            safe_name = re.sub(r'[<>:"/\\|?*]', '_', normalized_name)
            # 3. 연속된 공백을 하나로 변경
            safe_name = re.sub(r'\s+', '_', safe_name.strip())
            # 4. 빈 문자열이면 기본값 사용
            if not safe_name or safe_name == '_':
                safe_name = 'file'
            
            # 최종 파일명 생성
            filename = f"{safe_name}.{file_ext}"
            print(f"🔍 최종 안전한 파일명: '{filename}'")
            if file_ext not in ALLOWED_EXTENSIONS:
                return jsonify({
                    'success': False, 
                    'error': '지원되는 파일 형식: PDF, DOCX, JPG, JPEG, PNG, GIF, BMP',
                    'conversion_method': conversion_method
                }), 400
            
            # 파일 저장
            import time
            timestamp = str(int(time.time()))
            safe_filename = f"{timestamp}_{filename}"
            input_path = os.path.join(UPLOAD_FOLDER, safe_filename)
            
            try:
                os.makedirs(UPLOAD_FOLDER, exist_ok=True)
                file.save(input_path)
                
                saved_file_size = os.path.getsize(input_path)
                if saved_file_size == 0:
                    os.remove(input_path)
                    return jsonify({
                        'success': False, 
                        'error': '파일 저장 중 오류가 발생했습니다.',
                        'conversion_method': conversion_method
                    }), 500
                
                print(f"파일 저장 완료 - 크기: {saved_file_size}바이트")
            except Exception as e:
                print(f"파일 저장 오류: {str(e)}")
                return jsonify({
                    'success': False, 
                    'error': f'파일 저장 중 오류가 발생했습니다: {str(e)}',
                    'conversion_method': conversion_method
                }), 500
            
            # 변환 처리
            conversion_success = False
            output_path = None
            
            try:
                if file_ext == 'pdf':
                    # PDF → DOCX 변환 - 출력 파일명 미리 고정
                    base = filename.rsplit(".", 1)[0] if "." in filename else filename
                    output_filename = base + ".docx"
                    output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                    
                    quality = request.form.get('quality', 'medium')
                    print(f"PDF → DOCX 변환 시작 - {input_path} -> {output_path}")
                    
                    # 암호화된 PDF 체크
                    import fitz
                    try:
                        with fitz.open(input_path) as doc:
                            if doc.isEncrypted:
                                return jsonify(success=False, error="ENCRYPTED_PDF"), 400
                    except Exception as e:
                        print(f"PDF 암호화 체크 실패: {e}")
                    
                    # ADOBE_DISABLED 환경변수 체크
                    if os.getenv("ADOBE_DISABLED") == "true":
                        adobe_ready = False
                    
                    if adobe_ready:
                        print(">>> [DEBUG] image-only or vector PDF detected -> try Adobe first", flush=True)
                        ok, info = adobe_pdf_to_docx(input_path, output_path)
                        if not ok:
                            print(">>> [DEBUG] Adobe failed -> fallback to pdf2docx/image_to_docx", flush=True)
                            ok = pdf_to_docx(input_path, output_path, quality)
                            if not ok:
                                return jsonify(success=False, error="ADOBE_AND_FALLBACK_FAILED", detail=info), 400
                    else:
                        ok = pdf_to_docx(input_path, output_path, quality)
                        if not ok:
                            return jsonify(success=False, error="PDF2DOCX_FAIL"), 400
                    
                    conversion_success = True
                    
                elif file_ext == 'docx':
                    # DOCX → PDF 변환
                    base_filename = filename.rsplit('.', 1)[0] if '.' in filename else filename
                    output_filename = base_filename + '.pdf'
                    output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                    
                    print(f"DOCX → PDF 변환 시작 - {input_path} -> {output_path}")
                    conversion_success = docx_to_pdf(input_path, output_path)
                    
                elif file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                    # 이미지 → DOCX 변환
                    base_filename = filename.rsplit('.', 1)[0] if '.' in filename else filename
                    output_filename = base_filename + '.docx'
                    output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                    
                    print(f"이미지 → DOCX 변환 시작 - {input_path} -> {output_path}")
                    conversion_success = image_to_docx(input_path, output_path)
                
                # 변환 결과 처리
                if conversion_success:
                    print("변환 성공 - 다운로드 준비")
                    
                    # 임시 파일 정리
                    try:
                        os.remove(input_path)
                        print("임시 파일 삭제 완료")
                    except Exception as e:
                        print(f"임시 파일 삭제 실패 (무시됨): {e}")
                    
                    # 파일 다운로드 제공
                    try:
                        print("파일 다운로드 시작")
                        return send_file(output_path, as_attachment=True, download_name=output_filename)
                    except Exception as e:
                        print(f"파일 다운로드 오류: {str(e)}")
                        return jsonify({
                            'success': False, 
                            'error': f'파일 다운로드 중 오류가 발생했습니다: {str(e)}',
                            'conversion_method': conversion_method
                        }), 500
                else:
                    print("변환 실패 - 정리 작업")
                    
                    # 실패한 파일들 정리
                    for cleanup_path in [input_path, output_path]:
                        try:
                            if cleanup_path and os.path.exists(cleanup_path):
                                os.remove(cleanup_path)
                        except Exception as e:
                            print(f"파일 정리 실패 (무시됨): {e}")
                    
                    return jsonify({
                        'success': False, 
                        'error': '파일 변환에 실패했습니다.',
                        'conversion_method': conversion_method
                    }), 500
                    
            except Exception as e:
                print(f"변환 중 예외 발생: {str(e)}")
                return jsonify({
                    'success': False, 
                    'error': f'변환 중 오류가 발생했습니다: {str(e)}',
                    'conversion_method': conversion_method
                }), 500
        else:
            return jsonify({
                'success': False, 
                'error': '지원되지 않는 파일 형식입니다.',
                'conversion_method': conversion_method
            }), 400
            
    except Exception as e:
        print(f"업로드 처리 중 예외 발생: {str(e)}")
        return jsonify({
            'success': False, 
            'error': f'파일 처리 중 오류가 발생했습니다: {str(e)}'
        }), 500

@app.route('/upload', methods=['POST'])
def upload_file():
    """기존 웹 인터페이스용 업로드 (리다이렉트 방식)"""
    try:
        print("파일 업로드 요청 시작")
        
        # 1단계: 파일 존재 여부 확인
        if 'file' not in request.files:
            flash('파일이 선택되지 않았습니다.')
            return redirect(url_for('index'))
        
        file = request.files['file']
        
        # 2단계: 파일명 확인
        if not file or file.filename == '' or file.filename is None:
            flash('파일이 선택되지 않았습니다.')
            return redirect(url_for('index'))
        
        # 3단계: 파일 내용 및 크기 확인 (강화된 검증)
        try:
            # 파일 포인터를 끝으로 이동하여 크기 확인
            file.seek(0, 2)
            file_size = file.tell()
            file.seek(0)  # 파일 포인터를 다시 처음으로 이동
            
            # 파일 크기가 0인 경우 처리
            if file_size == 0:
                flash('업로드된 파일이 비어있습니다. 올바른 PDF 파일을 선택해주세요.')
                return redirect(url_for('index'))
            
            # 최소 파일 크기 확인 (PDF 헤더 최소 크기)
            if file_size < 100:  # 100바이트 미만은 유효한 PDF가 아님
                flash('파일이 너무 작습니다. 올바른 PDF 파일을 선택해주세요.')
                return redirect(url_for('index'))
            
            if file_size > 100 * 1024 * 1024:  # 100MB
                flash(f'파일 크기가 너무 큽니다. (현재: {file_size // (1024*1024)}MB, 최대: 100MB)')
                return redirect(url_for('index'))
            
            print(f"파일 크기: {file_size // (1024*1024) if file_size >= 1024*1024 else file_size // 1024}{'MB' if file_size >= 1024*1024 else 'KB'}")
            
            # 4단계: 파일 형식 검증
            file_content = file.read(10)  # 처음 10바이트 읽기
            file.seek(0)  # 다시 처음으로 이동
            
            # 파일 확장자 확인 (filename을 file.filename으로 수정)
            file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
            
            # PDF 파일 검증
            if file_ext == 'pdf' and not file_content.startswith(b'%PDF-'):
                flash('올바른 PDF 파일이 아닙니다. PDF 파일을 선택해주세요.')
                return redirect(url_for('index'))
            
            # 이미지 파일 검증 (기본적인 헤더 체크)
            elif file_ext in ['jpg', 'jpeg'] and not (file_content.startswith(b'\xff\xd8\xff') or file_content.startswith(b'\xff\xd8')):
                flash('올바른 JPEG 파일이 아닙니다.')
                return redirect(url_for('index'))
            
            elif file_ext == 'png' and not file_content.startswith(b'\x89PNG\r\n\x1a\n'):
                flash('올바른 PNG 파일이 아닙니다.')
                return redirect(url_for('index'))
            
            elif file_ext == 'gif' and not (file_content.startswith(b'GIF87a') or file_content.startswith(b'GIF89a')):
                flash('올바른 GIF 파일이 아닙니다.')
                return redirect(url_for('index'))
            
            elif file_ext == 'bmp' and not file_content.startswith(b'BM'):
                flash('올바른 BMP 파일이 아닙니다.')
                return redirect(url_for('index'))
                
        except Exception as e:
            print(f"파일 검증 중 오류: {str(e)}")
            flash('파일을 읽는 중 오류가 발생했습니다. 다른 파일을 시도해주세요.')
            return redirect(url_for('index'))
        
        # 5단계: 파일 형식 확인 및 처리 (강화된 검증)
        if file and allowed_file(file.filename):
            # 원본 파일명에서 확장자 추출
            original_filename = file.filename
            print(f"🔍 [upload] 원본 파일명: '{original_filename}'")
            
            # 확장자 추출 (원본 파일명에서)
            if '.' not in original_filename:
                flash('파일 확장자가 없습니다. PDF 파일을 선택해주세요.')
                return redirect(url_for('index'))
            
            file_ext = original_filename.rsplit('.', 1)[1].lower()
            print(f"🔍 [upload] 추출된 확장자: '{file_ext}'")
            
            # 안전한 파일명 생성 (한글 파일명 처리)
            import re
            import unicodedata
            
            # 원본 파일명에서 확장자 제거
            base_name = original_filename.rsplit('.', 1)[0]
            
            # 한글 파일명을 안전하게 처리
            # 1. 유니코드 정규화
            normalized_name = unicodedata.normalize('NFC', base_name)
            # 2. 안전하지 않은 문자 제거 (경로 구분자, 특수문자 등)
            safe_name = re.sub(r'[<>:"/\\|?*]', '_', normalized_name)
            # 3. 연속된 공백을 하나로 변경
            safe_name = re.sub(r'\s+', '_', safe_name.strip())
            # 4. 빈 문자열이면 기본값 사용
            if not safe_name or safe_name == '_':
                safe_name = 'file'
            
            # 최종 파일명 생성
            filename = f"{safe_name}.{file_ext}"
            print(f"🔍 [upload] 최종 안전한 파일명: '{filename}'")
            if file_ext not in ALLOWED_EXTENSIONS:
                flash('지원되는 파일 형식: PDF, DOCX, JPG, JPEG, PNG, GIF, BMP')
                return redirect(url_for('index'))
            
            # 안전한 파일명 생성 (타임스탬프 추가로 중복 방지)
            import time
            timestamp = str(int(time.time()))
            safe_filename = f"{timestamp}_{filename}"
            input_path = os.path.join(UPLOAD_FOLDER, safe_filename)
            
            print(f"파일 저장 중 - {input_path}")
            try:
                # 파일 저장 전 디렉토리 존재 확인
                os.makedirs(UPLOAD_FOLDER, exist_ok=True)
                file.save(input_path)
                
                # 저장된 파일 크기 재확인
                saved_file_size = os.path.getsize(input_path)
                if saved_file_size == 0:
                    os.remove(input_path)
                    flash('파일 저장 중 오류가 발생했습니다. 다시 시도해주세요.')
                    return redirect(url_for('index'))
                
                print(f"파일 저장 완료 - 크기: {saved_file_size}바이트")
            except Exception as e:
                print(f"파일 저장 오류: {str(e)}")
                flash(f'파일 저장 중 오류가 발생했습니다: {str(e)}')
                return redirect(url_for('index'))
            
            # 변환 처리
            conversion_success = False
            output_path = None
            
            if file_ext == 'pdf':
                # PDF → DOCX 변환
                # 파일명에서 확장자 제거 (안전하게)
                base_filename = filename.rsplit('.', 1)[0] if '.' in filename else filename
                output_filename = base_filename + '.docx'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                quality = request.form.get('quality', 'medium')
                print(f"PDF → DOCX 변환 시작 - {input_path} -> {output_path}")
                
                try:
                    conversion_result = pdf_to_docx(input_path, output_path, quality)
                    
                    # --- 핵심 수정 부분: OCR 텍스트 추출 실패 처리 ---
                    if conversion_result is None:
                        # 변환 실패 (텍스트 없음)
                        # 임시 파일 정리
                        try:
                            os.remove(input_path)
                        except:
                            pass
                        return jsonify({
                            'success': False, 
                            'error': 'PDF 파일에서 텍스트를 추출할 수 없습니다. 이미지 품질이 낮거나 텍스트가 없는 파일일 수 있습니다.',
                            'conversion_method': conversion_method
                        }), 400
                    else:
                        conversion_success = conversion_result
                except Exception as e:
                    print(f"변환 중 예외 발생: {str(e)}")
                    flash(f'변환 중 오류가 발생했습니다: {str(e)}')
                    
            elif file_ext == 'docx':
                # DOCX → PDF 변환
                # 파일명에서 확장자 제거 (안전하게)
                base_filename = filename.rsplit('.', 1)[0] if '.' in filename else filename
                output_filename = base_filename + '.pdf'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                print(f"DOCX → PDF 변환 시작 - {input_path} -> {output_path}")
                
                try:
                    conversion_success = docx_to_pdf(input_path, output_path)
                except Exception as e:
                    print(f"변환 중 예외 발생: {str(e)}")
                    flash(f'변환 중 오류가 발생했습니다: {str(e)}')
                    
            elif file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                # 이미지 → DOCX 변환
                # 파일명에서 확장자 제거 (안전하게)
                base_filename = filename.rsplit('.', 1)[0] if '.' in filename else filename
                output_filename = base_filename + '.docx'
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                print(f"이미지 → DOCX 변환 시작 - {input_path} -> {output_path}")
                
                try:
                    conversion_success = image_to_docx(input_path, output_path)
                except Exception as e:
                    print(f"변환 중 예외 발생: {str(e)}")
                    flash(f'변환 중 오류가 발생했습니다: {str(e)}')
            
            # 변환 결과 처리
            if conversion_success:
                print("변환 성공 - 다운로드 준비")
                
                # 업로드된 파일 정리
                try:
                    os.remove(input_path)
                    print("임시 파일 삭제 완료")
                except Exception as e:
                    print(f"임시 파일 삭제 실패 (무시됨): {e}")
                
                # 파일 다운로드 제공
                try:
                    print("파일 다운로드 시작")
                    return send_file(output_path, as_attachment=True, download_name=output_filename)
                except Exception as e:
                    print(f"파일 다운로드 오류: {str(e)}")
                    flash(f'파일 다운로드 중 오류가 발생했습니다: {str(e)}')
                    return redirect(url_for('index'))
            else:
                print("변환 실패 - 정리 작업")
                flash('파일 변환에 실패했습니다. 다시 시도해주세요.')
                
                # 실패한 파일들 정리
                for cleanup_path in [input_path, output_path]:
                    try:
                        if cleanup_path and os.path.exists(cleanup_path):
                            os.remove(cleanup_path)
                    except Exception as e:
                        print(f"파일 정리 실패 (무시됨): {e}")
                
                return redirect(url_for('index'))
        else:
            flash('PDF 또는 DOCX 파일만 업로드 가능합니다.')
            return redirect(url_for('index'))
            
    except Exception as e:
        print(f"업로드 처리 중 예외 발생: {str(e)}")
        flash('파일 처리 중 오류가 발생했습니다.')
        return redirect(url_for('index'))

if __name__ == '__main__':
    port = int(os.environ.get("PORT", "5000"))
    app.run(debug=True, host='0.0.0.0', port=port)