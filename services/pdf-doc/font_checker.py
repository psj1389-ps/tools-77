import os
import platform
from pathlib import Path

def check_system_fonts():
    """
    시스템에 설치된 한글 폰트 확인
    """
    print("=== 시스템 폰트 확인 ===")
    
    # Windows 폰트 경로
    if platform.system() == "Windows":
        font_paths = [
            "C:\\Windows\\Fonts",
            os.path.expanduser("~\\AppData\\Local\\Microsoft\\Windows\\Fonts")
        ]
    else:
        font_paths = ["/usr/share/fonts", "/System/Library/Fonts"]
    
    korean_fonts = []
    target_fonts = ['malgun', 'noto', 'gulim', 'dotum', 'batang']
    
    for font_path in font_paths:
        if os.path.exists(font_path):
            for font_file in os.listdir(font_path):
                font_lower = font_file.lower()
                if any(target in font_lower for target in target_fonts):
                    korean_fonts.append(font_file)
    
    print(f"발견된 한글 폰트: {len(korean_fonts)}개")
    for font in korean_fonts[:10]:  # 처음 10개만 표시
        print(f"  - {font}")
    
    # 맑은 고딕 확인
    malgun_fonts = [f for f in korean_fonts if 'malgun' in f.lower()]
    if malgun_fonts:
        print(f"\n✅ 맑은 고딕 폰트 발견: {malgun_fonts}")
        return True, '맑은 고딕'
    
    # Noto Sans CJK 확인
    noto_fonts = [f for f in korean_fonts if 'noto' in f.lower() and 'cjk' in f.lower()]
    if noto_fonts:
        print(f"\n✅ Noto Sans CJK 폰트 발견: {noto_fonts}")
        return True, 'Noto Sans CJK KR'
    
    print("\n❌ 적절한 한글 폰트를 찾을 수 없습니다.")
    return False, None

def test_pptx_font():
    """
    PPTX에서 폰트 테스트
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # 테스트 프레젠테이션 생성
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 텍스트 박스 추가
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        text_frame = textbox.text_frame
        
        # 한글 텍스트 추가
        text_frame.text = "한글 폰트 테스트 - 맑은 고딕"
        paragraph = text_frame.paragraphs[0]
        
        # 폰트 설정 테스트
        fonts_to_test = ['맑은 고딕', 'Noto Sans CJK KR', 'Malgun Gothic', 'Gulim']
        
        for font_name in fonts_to_test:
            try:
                paragraph.font.name = font_name
                paragraph.font.size = Pt(14)
                print(f"✅ {font_name} 폰트 설정 성공")
                
                # 테스트 파일 저장
                test_file = f"font_test_{font_name.replace(' ', '_')}.pptx"
                prs.save(test_file)
                print(f"   테스트 파일 저장: {test_file}")
                
                # 파일 삭제
                os.remove(test_file)
                return font_name
                
            except Exception as e:
                print(f"❌ {font_name} 폰트 설정 실패: {e}")
        
        return None
        
    except Exception as e:
        print(f"PPTX 폰트 테스트 오류: {e}")
        return None

if __name__ == "__main__":
    # 시스템 폰트 확인
    font_available, recommended_font = check_system_fonts()
    
    # PPTX 폰트 테스트
    print("\n=== PPTX 폰트 테스트 ===")
    working_font = test_pptx_font()
    
    if working_font:
        print(f"\n🎉 권장 폰트: {working_font}")
    else:
        print("\n⚠️ 한글 폰트 설치가 필요합니다.")